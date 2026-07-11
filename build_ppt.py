from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper to create a slide with title and bullet points

def add_bullets(title, bullets, layout_index=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(22)
    return slide

slides = [
    ("CAB BOOKING SYSTEM", ["A Web Application Project", "By: Haritha"]),
    ("Objective", [
        "Provide easy online cab booking",
        "Calculate fare automatically",
        "Connect user with driver in real-time",
        "Save time and reduce manual work"
    ]),
    ("Features", [
        "User Login - Name and Phone Number",
        "Book Cab - Pickup and Drop location",
        "Payment Options - Cash, UPI, Card",
        "Driver Dashboard - See all bookings",
        "Fare Calculation - Based on distance"
    ]),
    ("Technologies Used", [
        "Frontend: HTML, CSS, JavaScript",
        "Backend: Node.js, Express.js",
        "Storage: LocalStorage / JSON file",
        "Deployment: GitHub + Render"
    ]),
    ("Working Process", [
        "User enters details and books cab",
        "Data is saved in server",
        "Driver can see booking on dashboard",
        "Driver accepts and completes ride"
    ]),
    ("Advantages", [
        "Fast and Easy Booking",
        "Mobile Friendly",
        "Transparent Fare",
        "24/7 Available"
    ]),
    ("Screenshots", [
        "home.png",
        "booking.png",
        "driver.png"
    ]),
    ("Conclusion", [
        "Cab Booking System makes travel easy and fast.",
        "It reduces waiting time and provides a smooth experience for both user and driver."
    ]),
]

for title, bullets in slides:
    add_bullets(title, bullets)

output_path = "cab_booking_presentation.pptx"
prs.save(output_path)
print(f"Created {output_path}")
