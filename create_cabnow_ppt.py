from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

root = r"C:\Users\haris\OneDrive\Desktop\cab booking"
os.chdir(root)
output = os.path.join(root, "CabNow_Presentation.pptx")

if os.path.exists(output):
    try:
        os.remove(output)
    except PermissionError:
        print("Existing PPT is locked; trying a temporary replacement path")
        output = os.path.join(root, "CabNow_Presentation_tmp.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

YELLOW = RGBColor(255, 215, 0)
BLACK = RGBColor(15, 15, 16)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(35, 35, 38)


def add_slide(title, lines, placeholder_boxes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.alignment = PP_ALIGN.CENTER
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = YELLOW

    body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.6))
    tf_body = body_box.text_frame
    tf_body.clear()
    tf_body.word_wrap = True
    tf_body.vertical_anchor = 1

    for i, line in enumerate(lines):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
        if line.startswith('•'):
            p.bullet = True
            p.level = 0

    if placeholder_boxes:
        x_positions = [1.0, 3.9, 6.8, 9.7]
        for idx, label in enumerate(placeholder_boxes):
            x = Inches(x_positions[idx])
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(4.5), Inches(2.2), Inches(1.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = GRAY
            shape.line.color.rgb = YELLOW
            tf_box = shape.text_frame
            tf_box.clear()
            p_box = tf_box.paragraphs[0]
            p_box.text = label
            p_box.alignment = PP_ALIGN.CENTER
            p_box.font.size = Pt(16)
            p_box.font.color.rgb = WHITE


add_slide('TITLE SLIDE', [
    'CabNow - Advanced Cab Booking App',
    'A Responsive Web Application',
    'Tech: HTML, CSS, Bootstrap 5, JavaScript',
    'Name: [Your Name]'
])

add_slide('PROBLEM STATEMENT', [
    'People face issues while booking cabs:',
    '• No quick car selection',
    '• No live tracking',
    '• No safety features',
    '• No discount options',
    'Solution: One app for all cab booking needs'
])

add_slide('KEY FEATURES', [
    '• 8 Car Categories: Mini to Luxury',
    '• Live Tracking + Driver Profile + SOS',
    '• Dark/Light Mode Toggle',
    '• Promo Code CAB20 - 20% OFF',
    '• Payment, Trip History, Rating',
    '• Fully Mobile Responsive'
])

add_slide('TECH STACK', [
    '• Frontend: HTML5, CSS3',
    '• Framework: Bootstrap 5',
    '• Language: JavaScript',
    '• Theme: Yellow #FFD700 + Black',
    '• Server: Node.js'
])

add_slide('APP SCREENSHOTS', ['Login Page | Car Selection | Live Tracking | Payment'], [
    'Login Page', 'Car Selection', 'Live Tracking', 'Payment'
])

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BLACK

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
tf_title = title_box.text_frame
tf_title.clear()
p_title = tf_title.paragraphs[0]
p_title.text = 'PROJECT STRUCTURE'
p_title.alignment = PP_ALIGN.CENTER
p_title.font.size = Pt(44)
p_title.font.bold = True
p_title.font.color.rgb = YELLOW

body_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.9), Inches(4.6))
tf_body = body_box.text_frame
tf_body.clear()
tf_body.word_wrap = True
lines = [
    'cab-booking-project/',
    '|__ index.html',
    '|__ booking.html, payment.html, driver.html',
    '|__ assets/',
    '|__ README.md',
    '|__ CabNow_Presentation.pptx'
]
for i, line in enumerate(lines):
    p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
    p.text = line
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.space_after = Pt(6)

add_slide('FUTURE ENHANCEMENTS', [
    '• Google Maps GPS Integration',
    '• Wallet and UPI Payment',
    '• Admin Panel for Drivers',
    '• Push Notifications'
])

add_slide('THANK YOU', ['Questions'])

prs.save(output)
print(f'Created {output}')
